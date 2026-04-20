"""
Generate a single-slide 1-minute presentation PowerPoint for ECE 510 project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# Background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

# ============================================================
# Title bar
# ============================================================
title_box = slide.shapes.add_shape(
    1, Inches(0), Inches(0), prs.slide_width, Inches(0.85))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
title_box.line.fill.background()

tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Accelerating Transformer ff_backward with a 64x64 Systolic Array Chiplet"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Subtitle
p2 = tf.add_paragraph()
p2.text = "ECE 510 - HW4AI  |  Spring 2026  |  David Boss"
p2.font.size = Pt(13)
p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p2.alignment = PP_ALIGN.CENTER

# ============================================================
# Helper: add a text box with title + bullet content
# ============================================================
def add_section(left, top, width, height, title, bullets, title_color=RGBColor(0x1B,0x3A,0x5C)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(4)

    # Bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(2)
        p.level = 0

    return txBox

# ============================================================
# Column 1: Questions a, b, c
# ============================================================
col1_left = 0.3
col1_top = 1.0
col1_w = 4.0

add_section(col1_left, col1_top, col1_w, 1.8,
    "a. What are we trying to do?",
    [
        "Accelerate the dominant kernel in transformer",
        "training: ff_backward (feed-forward backward",
        "pass), which accounts for 32.5% of total runtime.",
        "Target: 10x speedup via a custom systolic array",
        "chiplet with high-bandwidth on-chip SRAM.",
    ])

add_section(col1_left, col1_top + 1.9, col1_w, 1.8,
    "b. How have others done it?",
    [
        "Current baseline: pure NumPy on i5-10500H CPU.",
        "GPUs (RTX 3050 Ti): higher BW but FP64 is 1/64",
        "of FP32 throughput on consumer NVIDIA GPUs.",
        "Both platforms are memory-bandwidth bound -",
        "the kernel's AI (5.43 FLOP/B) < ridge points.",
    ])

add_section(col1_left, col1_top + 3.8, col1_w, 1.8,
    "c. What are we doing differently?",
    [
        "64x64 systolic array @ 500 MHz = 4.096 TFLOP/s.",
        "256 GB/s on-chip SRAM (10x CPU DRAM BW).",
        "UCIe x16 die-to-die link (8 GB/s, 45% utilized).",
        "Dimensions matched to d_model=64 for zero tiling",
        "overhead. Bandwidth-first design for memory-bound",
        "workloads.",
    ])

# ============================================================
# Column 2: Questions d, e + key metrics
# ============================================================
col2_left = 4.6
col2_top = 1.0
col2_w = 4.0

add_section(col2_left, col2_top, col2_w, 2.4,
    "d. What have we accomplished?",
    [
        "Profiled transformer over 1,000 iterations.",
        "Identified ff_backward as dominant kernel (32.5%).",
        "Computed AI = 5.43 FLOP/B analytically (FP64).",
        "Built roofline models for CPU, GPU, accelerator.",
        "Implemented CUDA GEMM (naive + tiled) on GPU.",
        "Designed HW/SW partition with UCIe interface.",
        "Baseline: 35.3 ms/iter, 3.4 GFLOP/s, 14.7 MB.",
    ])

add_section(col2_left, col2_top + 2.5, col2_w, 2.0,
    "e. What remains to be done?",
    [
        "RTL design of 64x64 systolic array in Verilog.",
        "SRAM controller + data tiling logic.",
        "UCIe interface integration.",
        "Functional verification and synthesis.",
        "Measure actual speedup vs. software baseline",
        "at M4 milestone (target: >= 10x).",
    ])

# Key metrics box
metrics_box = slide.shapes.add_shape(
    1, Inches(col2_left), Inches(col2_top + 4.7), Inches(col2_w), Inches(1.4))
metrics_box.fill.solid()
metrics_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
metrics_box.line.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
metrics_box.line.width = Pt(1)

tf = metrics_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Metrics"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

for line in [
    "Dominant kernel: ff_backward (32.5% runtime)",
    "Arithmetic Intensity: 5.43 FLOP/B (memory-bound)",
    "CPU baseline: 139 GFLOP/s | Target: 1,390 GFLOP/s",
    "Projected speedup: 10x",
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(1)

# ============================================================
# Column 3: Visualizations (roofline images)
# ============================================================
col3_left = 8.9
col3_w = 4.2

# Roofline from cf02
cf02_img = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\codefest\cf02\profiling\roofline_project.pdf"
cf03_img = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\codefest\cf03\profiling\gemm_roofline.png"

# Add cf03 combined roofline (has both cf02 and cf03 data)
if os.path.exists(cf03_img):
    slide.shapes.add_picture(cf03_img, Inches(col3_left), Inches(1.0), Inches(col3_w), Inches(2.8))

    # Label
    lbl = slide.shapes.add_textbox(Inches(col3_left), Inches(3.8), Inches(col3_w), Inches(0.3))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = "Combined Roofline: CPU / Accelerator / GPU"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

# Architecture summary diagram (text-based)
arch_box = slide.shapes.add_shape(
    1, Inches(col3_left), Inches(4.2), Inches(col3_w), Inches(2.8))
arch_box.fill.solid()
arch_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
arch_box.line.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
arch_box.line.width = Pt(1)

tf = arch_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "System Architecture"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
p.alignment = PP_ALIGN.CENTER

for line in [
    "",
    "  Host CPU (i5-10500H)",
    "     |  UCIe x16 @ 4 GT/s (8 GB/s)",
    "     v",
    "  +------ Accelerator Chiplet ------+",
    "  |  64x64 Systolic Array @ 500 MHz |",
    "  |  4,096 FP64 MAC units           |",
    "  |  Peak: 4.096 TFLOP/s            |",
    "  |                                 |",
    "  |  SRAM Scratchpad (256 GB/s)     |",
    "  +---------------------------------+",
]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(9)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.LEFT

# ============================================================
# Save
# ============================================================
out_path = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\presentation\1min_presentation.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
