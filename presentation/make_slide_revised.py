"""
Revised 1-minute presentation. Addresses critique:
1. Clarify full GEMM acceleration via systolic tiling
2. Justify attention traffic reduction (FlashAttention-style)
3. Reframe around dataflow fusion, not GELU
4. Frame as chiplet building block with tiling strategy
5. Qualify performance claims as roofline estimates
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

# ============================================================
# Title bar
# ============================================================
title_box = slide.shapes.add_shape(
    1, Inches(0), Inches(0), prs.slide_width, Inches(0.75))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
title_box.line.fill.background()

tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Dataflow-Fused Transformer Training Accelerator Chiplet"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "End-to-end fusion of FFN + attention sublayers to eliminate intermediate SRAM traffic  |  ECE 510  |  Spring 2026  |  David Boss"
p2.font.size = Pt(11)
p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p2.alignment = PP_ALIGN.CENTER

# ============================================================
# Helpers
# ============================================================
def add_box(left, top, width, height, fill_color, border_color=None):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box

def add_textbox(left, top, width, height, lines, title=None, title_size=13,
                body_size=9, title_color=RGBColor(0x1B,0x3A,0x5C)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(3)
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(body_size)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(1)
    return txBox

# ============================================================
# Column 1: Core idea + Performance
# ============================================================
c1 = 0.2

add_textbox(c1, 0.85, 4.2, 1.2,
    [
        "End-to-end fusion of transformer sublayers (FFN forward,",
        "FFN backward, multi-head attention) to eliminate intermediate",
        "SRAM traffic. All GEMMs execute on a 64x64 systolic array",
        "with tiled dataflow; activations (GELU, softmax) are fused",
        "at the array output — intermediates stay in registers.",
        "Targets ~80% of training runtime (B=4, T=64, D=64, 2 layers).",
    ],
    title="Core Idea: Dataflow Fusion Eliminates Memory Traffic")

# Performance table
tbl_box = add_box(c1, 2.2, 4.2, 2.5, RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1B,0x3A,0x5C))
tf = tbl_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FP32 Roofline Estimates (ff_backward kernel)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
p.alignment = PP_ALIGN.CENTER

rows = [
    "Platform           Peak      BW        Attain.   Bound",
    "──────────────────────────────────────────────────────────",
    "i5-10500H CPU      864G      25.6 GB/s    278 G/s   mem",
    "RTX 3050 Ti       5.3T     192 GB/s   2,085 G/s   mem",
    "RTX 4080         48.7T     717 GB/s   7,783 G/s   mem",
    "Accel+ (fused)    8.2T     512 GB/s   8,192 G/s   COMP",
    "",
    "Fusion shifts AI from 10.86 → 35.4 FLOP/B, crossing the",
    "ridge point. Accelerator becomes compute-bound, hitting",
    "the full 8.2 TFLOP/s systolic peak.",
    "",
    "Note: estimates from roofline model, not measured silicon.",
    "Config: B=4, T=64, d_model=64, d_ff=256, 2 layers, FP32.",
]
for row in rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)

# Power
pwr_box = add_box(c1, 4.85, 4.2, 1.5, RGBColor(0xE8,0xF0,0xFE), RGBColor(0x1B,0x3A,0x5C))
tf = pwr_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Power Efficiency (estimated)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

pwr_rows = [
    "Platform         Power    Perf       Eff. (est.)",
    "─────────────────────────────────────────────────",
    "i5-10500H CPU    45W      278 G/s      6.2 G/s/W",
    "RTX 4080        320W    7,783 G/s     24.3 G/s/W",
    "Accel+ (est.)  ~11.5W   8,192 G/s   ~712  G/s/W",
    "",
    "Power est. from component analysis, not synthesis.",
]
for row in pwr_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)

# ============================================================
# Column 2: Dataflow + Architecture + Future
# ============================================================
c2 = 4.7

# Dataflow description — reframed around dataflow, not GELU
hw_box = add_box(c2, 0.85, 4.2, 2.8, RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1B,0x3A,0x5C))
tf = hw_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Accelerator Dataflow"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

hw_rows = [
    "Full FFN + attention pipeline on systolic array via tiling:",
    "  GEMMs tiled as 64x64 blocks through the array (D=64",
    "  maps directly; F=256 requires 4 passes per tile dim).",
    "",
    "Fused activation units at systolic output stage:",
    "  MUX selects GELU, GELU', or softmax per FSM state.",
    "  Intermediates flow through registers, never to SRAM.",
    "",
    "Attention uses FlashAttention-style streaming/tiling:",
    "  Q*K^T computed in blocks, softmax applied per-block,",
    "  avoids materializing the full T*T score matrix.",
    "",
    "Memory traffic reduction (vs unfused):",
    "  ff_backward:  3.2 MB → 0.98 MB  (69% reduction)",
    "  attention:    1.9 MB → 0.33 MB  (83% reduction)",
]
for row in hw_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)

# Datapath diagram — tile flow, fusion points, memory avoidance
dp_box = add_box(c2, 3.8, 4.2, 2.6, RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0x1B, 0x3A, 0x5C))
tf = dp_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Fused Datapath: Tile Flow Through Systolic Array"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
p.alignment = PP_ALIGN.CENTER

future_rows = [
    "",
    " SRAM                  Systolic Array              SRAM",
    " ─────                 ──────────────              ─────",
    "                    ┌──────────────────┐",
    " Tile A (64x64) ──>│                  │",
    "                    │  64x64 MAC grid  │──> partial",
    " Tile B (64x64) ──>│  (tiled GEMM)    │    sums",
    "                    └────────┬─────────┘",
    "                             │ (in registers, NOT to SRAM)",
    "                             v",
    "                    ┌──────────────────┐",
    "                    │ Fused Activation  │",
    "                    │ MUX: GELU|GELU'  │",
    "                    │      |softmax    │",
    "                    └────────┬─────────┘",
    "                             │ (still in registers)",
    "                             v",
    "                    ┌──────────────────┐",
    "                    │  Next tiled GEMM  │──> final",
    "                    │  (reuses array)   │    output ──> SRAM",
    "                    └──────────────────┘",
    "",
    " Memory avoided: intermediate tensors (dh_act, gelu_grad",
    " result, scores, attn) never written to / read from SRAM.",
]
for row in future_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x55, 0x44, 0x33)
    p.space_after = Pt(0)

# ============================================================
# Column 3: Roofline + Architecture diagram
# ============================================================
c3 = 9.15
c3_w = 4.0

roofline_img = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\codefest\cf03\profiling\all.png"
if os.path.exists(roofline_img):
    slide.shapes.add_picture(roofline_img, Inches(c3), Inches(0.85), Inches(c3_w), Inches(2.8))
    lbl = slide.shapes.add_textbox(Inches(c3), Inches(3.65), Inches(c3_w), Inches(0.25))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = "FP32 Roofline: ff_backward + attention on all platforms (roofline model)"
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

# Architecture diagram
arch_box = add_box(c3, 4.0, c3_w, 2.5, RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1B,0x3A,0x5C))
tf = arch_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Chiplet Architecture (building block)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
p.alignment = PP_ALIGN.CENTER

arch_rows = [
    "",
    "  Host CPU (i5-10500H)    ← 20% of workload",
    "     |  UCIe x16 (8 GB/s bidir)",
    "     v",
    "  +───── Accelerator+ Chiplet ──────+",
    "  │                                 │",
    "  │  64x64 Systolic Array @ 500 MHz │",
    "  │  4,096 MACs, tiled GEMM engine  │",
    "  │          |                      │",
    "  │   [MUX: GELU|GELU'|softmax]    │",
    "  │   Fused at output stage         │",
    "  │                                 │",
    "  │  SRAM Scratchpad (512 GB/s)     │",
    "  │  ~11.5W est.                    │",
    "  +─────────────────────────────────+",
]
for row in arch_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.LEFT

# Bottom bar
bot = slide.shapes.add_shape(
    1, Inches(0), Inches(7.15), prs.slide_width, Inches(0.35))
bot.fill.solid()
bot.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
bot.line.fill.background()
tf = bot.text_frame
p = tf.paragraphs[0]
p.text = "Dataflow fusion across FFN + attention sublayers  |  AI: 10.86 → 35.4 (mem→compute bound)  |  ~80% coverage, 3.6x Amdahl's  |  Chiplet building block, scalable via UCIe mesh"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Save
out_path = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\presentation\revised_v2.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
