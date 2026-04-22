"""
Generate a single-slide 1-minute presentation PowerPoint for ECE 510 project.
Accelerator+ version: includes attention, fusion, 512 GB/s SRAM.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

# ============================================================
# Title bar
# ============================================================
title_box = slide.shapes.add_shape(
    1, Inches(0), Inches(0), prs.slide_width, Inches(0.75))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
title_box.line.fill.background()

tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Accelerator+: Fused ff_backward + Attention Systolic Chiplet"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "ECE 510 - HW4AI  |  Spring 2026  |  David Boss"
p2.font.size = Pt(12)
p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p2.alignment = PP_ALIGN.CENTER

# ============================================================
# Helper
# ============================================================
def add_box(left, top, width, height, fill_color, border_color=None):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box

def add_textbox(left, top, width, height, lines, title=None, title_size=13,
                body_size=9, title_color=RGBColor(0x1B,0x3A,0x5C)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(3)
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(body_size)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(1)
    return txBox

# ============================================================
# Column 1: Problem + Performance Comparison
# ============================================================
c1 = 0.2

add_textbox(c1, 0.85, 4.2, 1.0,
    [
        "Accelerated: ff_backward (32.5%) + ff_forward (23.5%) + attention (24.8%)",
        "Combined: 80% of training runtime on hardware",
        "ff_forward reuses same systolic array + gelu circuit (zero extra area)",
        "Operator fusion eliminates intermediate memory traffic",
        "ff_backward fused AI: 10.86 → 35.4 | attn fused AI: 20.1 → 116.1",
    ],
    title="The Approach: Fuse & Accelerate 80% of Training")

# Performance comparison table
tbl_box = add_box(c1, 2.05, 4.2, 2.8, RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1B,0x3A,0x5C))
tf = tbl_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FP32 Performance Comparison (ff_backward)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
p.alignment = PP_ALIGN.CENTER

rows = [
    "Platform           Peak      BW        Attain.   Status",
    "──────────────────────────────────────────────────────────",
    "i5-10500H CPU      864G      25.6 GB/s    278 G/s   mem-bound",
    "RTX 3050 Ti GPU   5.3T     192 GB/s   2,085 G/s   mem-bound",
    "RTX 4080 GPU     48.7T     717 GB/s   7,783 G/s   mem-bound",
    "Accel+ (fused)    8.2T     512 GB/s   8,192 G/s   COMP-BOUND",
    "",
    "Fusion shifts AI from 10.86 → 35.4 FLOP/B, crossing the",
    "ridge point (16.0). Accel+ becomes COMPUTE-BOUND, hitting",
    "the full 8.2 TFLOP/s peak — beating even the RTX 4080.",
]

for row in rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)

# Power comparison
pwr_box = add_box(c1, 5.0, 4.2, 1.4, RGBColor(0xE8,0xF0,0xFE), RGBColor(0x1B,0x3A,0x5C))
tf = pwr_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Power & Efficiency"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

pwr_rows = [
    "Platform         Power    Perf       Efficiency",
    "─────────────────────────────────────────────────",
    "i5-10500H CPU    45W      278 G/s      6.2 GFLOP/s/W",
    "RTX 4080        320W    7,783 G/s     24.3 GFLOP/s/W",
    "Accel+         ~11.5W   8,192 G/s    712.3 GFLOP/s/W",
    "",
    "Accel+: beats 4080 at 1/28 power → 29x efficiency",
]
for row in pwr_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)

# ============================================================
# Column 2: HW Mapping + Future
# ============================================================
c2 = 4.7

hw_box = add_box(c2, 0.85, 4.2, 3.0, RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1B,0x3A,0x5C))
tf = hw_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Python → Hardware Translation"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

hw_rows = [
    "Python Concept         Hardware Equivalent",
    "──────────────────────────────────────────────",
    "@ (matmul)             64x64 systolic array (4096 MACs)",
    "gelu_grad()            Pipelined LUT + multipliers",
    "softmax()              Exp LUT + divider + comparators",
    "* (elementwise)        Single-cycle multiplier",
    ".sum() (reduction)     Adder tree",
    "Variable assignment    SRAM read/write (512 GB/s)",
    "Function call order    FSM state transitions (~25 states)",
    "",
    "FUSION: intermediates never touch memory.",
    "  systolic → [gelu_grad/softmax] → systolic",
    "  Results flow through registers, not SRAM.",
    "  Eliminates 69% of ff_bwd memory traffic",
    "  Eliminates 83% of attention memory traffic",
]
for row in hw_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)

# Future enhancements
future_box = add_box(c2, 4.0, 4.2, 2.4, RGBColor(0xFE, 0xF3, 0xE0), RGBColor(0xE6, 0x7E, 0x22))
tf = future_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Future Enhancements"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0xC0, 0x5C, 0x00)

future_rows = [
    "1. Mixed-precision (FP16/BF16 accumulate to FP32)",
    "   → 2x MACs per cycle, 2x effective bandwidth",
    "   → Peak jumps to 16.4 TFLOP/s at same area/power",
    "",
    "2. Accelerate layer_norm (~11% runtime)",
    "   → Coverage: 80% → ~91%, Amdahl's: 3.6x → 5.3x",
    "",
    "3. Double-buffered SRAM for overlapping compute + DMA",
    "   → Hides UCIe transfer latency completely",
    "",
    "4. Configurable array partitioning (2x 32x64)",
    "   → Run two matmuls concurrently for smaller dims",
]
for row in future_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x55, 0x44, 0x33)
    p.space_after = Pt(0)

# ============================================================
# Column 3: Roofline plot + Architecture
# ============================================================
c3 = 9.15
c3_w = 4.0

roofline_img = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\codefest\cf03\profiling\all.png"
if os.path.exists(roofline_img):
    slide.shapes.add_picture(roofline_img, Inches(c3), Inches(0.85), Inches(c3_w), Inches(2.8))
    lbl = slide.shapes.add_textbox(Inches(c3), Inches(3.65), Inches(c3_w), Inches(0.25))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = "FP32 Roofline: ff_backward + attention — all platforms"
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

# Architecture diagram
arch_box = add_box(c3, 4.0, c3_w, 2.5, RGBColor(0xFF,0xFF,0xFF), RGBColor(0x1B,0x3A,0x5C))
tf = arch_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Accelerator+ Architecture"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
p.alignment = PP_ALIGN.CENTER

arch_rows = [
    "",
    "  Host CPU (i5-10500H)    ← runs 20% of workload",
    "     |  UCIe x16 (8 GB/s bidir)",
    "     v",
    "  +──────── Accelerator+ Chiplet ────────+",
    "  │  64x64 Systolic Array @ 500 MHz      │",
    "  │  4,096 MACs → 8.2 TFLOP/s FP32      │",
    "  │                                      │",
    "  │  [gelu_grad] ← fused at output       │",
    "  │  [softmax]   ← fused at output       │",
    "  │  [causal mask] ← gate logic          │",
    "  │                                      │",
    "  │  SRAM Scratchpad (512 GB/s)  ~11.5W  │",
    "  +──────────────────────────────────────+",
]
for row in arch_rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(8.5)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.LEFT

# Bottom bar
bot = slide.shapes.add_shape(
    1, Inches(0), Inches(7.15), prs.slide_width, Inches(0.35))
bot.fill.solid()
bot.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
bot.line.fill.background()
tf = bot.text_frame
p = tf.paragraphs[0]
p.text = "80% of training accelerated  |  Fusion: AI 10.86 → 35.4 (compute-bound)  |  8,192 GFLOP/s @ 11.5W  |  29x more efficient than RTX 4080  |  Amdahl's: 3.6x end-to-end"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# Save
out_path = r"C:\Users\david\OneDrive\Documents\psu\510\GitHub\presentation\1min_presentation_accel+.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
